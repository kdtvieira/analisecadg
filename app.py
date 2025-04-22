import streamlit as st
import pandas as pd
import plotly.express as px
import plotly.graph_objects as go
import folium
from folium.plugins import HeatMap
from streamlit_folium import folium_static
import datetime
import io
from openpyxl import Workbook
from openpyxl.utils.dataframe import dataframe_to_rows
from openpyxl.chart import BarChart, Reference, PieChart
from openpyxl.drawing.image import Image
import matplotlib.pyplot as plt
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from geopy.geocoders import Nominatim
from geopy.exc import GeocoderTimedOut, GeocoderUnavailable
import time
import uuid
from PIL import Image as PILImage

# Configuração da página
st.set_page_config(
    page_title="Análise de Dados Criminais - PM",
    page_icon="🚨",
    layout="wide",
    initial_sidebar_state="collapsed"  # Iniciar com sidebar recolhida para maximizar espaço
)

# Definir estilos CSS personalizados para melhorar a visibilidade
st.markdown("""
<style>
    /* Aumentar tamanho dos títulos */
    .main h1 {
        font-size: 2.8rem !important;
        padding-bottom: 1rem;
    }
    .main h2 {
        font-size: 2.2rem !important;
        padding-bottom: 0.8rem;
    }
    .main h3 {
        font-size: 1.8rem !important;
        padding-bottom: 0.6rem;
    }
    
    /* Aumentar espaçamento e tamanho dos elementos */
    .stPlotlyChart {
        height: 600px !important;
    }
    
    /* Melhorar visibilidade dos widgets */
    .stSelectbox, .stMultiselect {
        margin-bottom: 1.5rem;
    }
    
    /* Melhorar visibilidade das métricas */
    .stMetric {
        background-color: #f0f2f6;
        padding: 1rem;
        border-radius: 0.5rem;
        box-shadow: 0 1px 2px rgba(0,0,0,0.1);
    }
    
    /* Estilizar abas para maior destaque */
    .stTabs [data-baseweb="tab-list"] {
        gap: 2rem;
    }
    .stTabs [data-baseweb="tab"] {
        height: 3rem;
        white-space: pre-wrap;
        font-size: 1rem;
        font-weight: 600;
        background-color: #f0f2f6;
        border-radius: 0.5rem 0.5rem 0 0;
    }
    .stTabs [aria-selected="true"] {
        background-color: #e0e5f0;
    }
    
    /* Estilo para o rodapé com créditos do autor */
    .footer {
        position: fixed;
        right: 15px;
        bottom: 10px;
        color: #666;
        font-size: 0.8rem;
        font-style: italic;
        z-index: 999;
    }
    
    /* Estilo para o logo no canto superior direito */
    .logo-container {
        position: fixed;
        top: 15px;
        right: 15px;
        z-index: 999;
    }
</style>

<!-- Adicionar créditos do autor no canto inferior direito -->
<div class="footer">
    Criado por Leandro Vieira de Souza
</div>
""", unsafe_allow_html=True)

# Adicionar logo no canto superior direito
logo_path = "assets/pmms_logo.png"
logo_image = PILImage.open(logo_path)
logo_col1, logo_col2 = st.columns([4, 1])
with logo_col2:
    st.image(logo_image, width=150)

# Inicializar estado da sessão para armazenar múltiplas planilhas
if 'dataframes' not in st.session_state:
    st.session_state.dataframes = {}  # Dicionário para armazenar múltiplos DataFrames

if 'active_dataframes' not in st.session_state:
    st.session_state.active_dataframes = []  # Lista para controlar quais DataFrames estão ativos

# Meses do ano para seleção
MESES = [
    "Janeiro", "Fevereiro", "Março", "Abril", "Maio", "Junho",
    "Julho", "Agosto", "Setembro", "Outubro", "Novembro", "Dezembro"
]

# Função para extrair unidades individuais de uma string com múltiplas unidades
def extract_units(unit_string):
    if pd.isna(unit_string):
        return []
    
    # Dividir por ponto e vírgula para separar múltiplas unidades
    units = [unit.strip() for unit in str(unit_string).split(';')]
    return units

# Função para carregar os dados
@st.cache_data
def load_data(file, month_name=None):
    df = pd.read_excel(file)
    
    # Converter colunas de data e hora para datetime
    df['DATA_HORA'] = pd.to_datetime(
        df['DATA DE INÍCIO DO ATENDIMENTO'] + ' ' + df['HORA DE INÍCIO DO ATENDIMENTO'],
        format='%d/%m/%Y %H:%M:%S',
        errors='coerce'
    )
    
    # Adicionar coluna com o nome do mês para identificação
    if month_name:
        df['MES_REFERENCIA'] = month_name
    
    return df

# Função para combinar múltiplos DataFrames
def combine_dataframes(dataframes_dict, active_keys=None):
    if not dataframes_dict:
        return pd.DataFrame()
    
    # Se active_keys não for fornecido, use todas as chaves
    if active_keys is None or len(active_keys) == 0:
        active_keys = list(dataframes_dict.keys())
    
    # Filtrar apenas os DataFrames ativos
    active_dfs = [dataframes_dict[key] for key in active_keys if key in dataframes_dict]
    
    if not active_dfs:
        return pd.DataFrame()
    
    # Combinar os DataFrames
    combined_df = pd.concat(active_dfs, ignore_index=True)
    return combined_df

# Função para filtrar os dados
def filter_data(df, start_date, end_date, crime_type, location, unit, keywords):
    filtered_df = df.copy()
    
    # Filtro de data
    if start_date and end_date:
        filtered_df = filtered_df[
            (filtered_df['DATA_HORA'] >= pd.to_datetime(start_date)) & 
            (filtered_df['DATA_HORA'] <= pd.to_datetime(end_date))
        ]
    
    # Filtro de tipo de crime
    if crime_type:
        filtered_df = filtered_df[filtered_df['EVENTO'].isin(crime_type)]
    
    # Filtro de localidade
    if location:
        filtered_df = filtered_df[filtered_df['ÁREA URBANA'].isin(location)]
    
    # Filtro de unidade responsável - modificado para tratar múltiplas unidades
    if unit:
        # Criar uma máscara para filtrar registros que contêm qualquer uma das unidades selecionadas
        mask = filtered_df['UNIDADE DA VIATURA'].apply(
            lambda x: any(selected_unit in extract_units(x) for selected_unit in unit)
        )
        filtered_df = filtered_df[mask]
    
    # Filtro de palavras-chave
    if keywords:
        # Combinar históricos e evoluções para busca
        filtered_df = filtered_df[
            filtered_df['HISTÓRICOS'].str.contains(keywords, case=False, na=False) | 
            filtered_df['EVOLUÇÕES'].str.contains(keywords, case=False, na=False)
        ]
    
    return filtered_df

# Função para criar gráfico de barras
def create_bar_chart(df, column, title, color='#1E3A8A'):
    if df.empty:
        st.warning("Não há dados para exibir no gráfico.")
        return None
    
    count_df = df[column].value_counts().reset_index()
    count_df.columns = [column, 'Contagem']
    
    fig = px.bar(
        count_df, 
        x=column, 
        y='Contagem',
        title=title,
        color_discrete_sequence=[color],
        height=600  # Aumentar altura do gráfico
    )
    
    fig.update_layout(
        xaxis_title=column,
        yaxis_title="Número de Ocorrências",
        template="plotly_white",
        title_font_size=24,  # Aumentar tamanho do título
        font=dict(size=16),  # Aumentar tamanho da fonte geral
        margin=dict(l=50, r=50, t=80, b=50)  # Ajustar margens
    )
    
    return fig

# Função para criar gráfico de barras comparativo por mês
def create_comparative_bar_chart(df, column):
    if df.empty or 'MES_REFERENCIA' not in df.columns:
        st.warning("Não há dados para comparação entre meses.")
        return None
    
    # Agrupar por mês de referência e coluna selecionada
    grouped = df.groupby(['MES_REFERENCIA', column]).size().reset_index(name='Contagem')
    
    fig = px.bar(
        grouped,
        x=column,
        y='Contagem',
        color='MES_REFERENCIA',
        title=f"Comparação de {column} por Mês",
        barmode='group',
        height=600  # Aumentar altura do gráfico
    )
    
    fig.update_layout(
        xaxis_title=column,
        yaxis_title="Número de Ocorrências",
        legend_title="Mês",
        template="plotly_white",
        title_font_size=24,  # Aumentar tamanho do título
        font=dict(size=16),  # Aumentar tamanho da fonte geral
        legend=dict(font=dict(size=14)),  # Aumentar tamanho da fonte da legenda
        margin=dict(l=50, r=50, t=80, b=50)  # Ajustar margens
    )
    
    return fig

# Função para criar gráfico de variação percentual
def create_percentage_change_chart(df, column, months):
    if df.empty or 'MES_REFERENCIA' not in df.columns or len(months) < 2:
        st.warning("São necessários pelo menos dois meses para análise de variação.")
        return None
    
    # Agrupar por mês e coluna selecionada
    grouped = df.groupby(['MES_REFERENCIA', column]).size().reset_index(name='Contagem')
    
    # Criar um pivot para facilitar o cálculo
    pivot = grouped.pivot(index=column, columns='MES_REFERENCIA', values='Contagem').fillna(0)
    
    # Calcular a variação percentual entre os meses selecionados
    month1, month2 = months[0], months[1]
    
    if month1 not in pivot.columns or month2 not in pivot.columns:
        st.warning(f"Dados insuficientes para os meses {month1} e {month2}.")
        return None
    
    # Calcular variação percentual
    pivot['Variação'] = ((pivot[month2] - pivot[month1]) / pivot[month1] * 100).fillna(0)
    
    # Filtrar apenas os tipos de crime com dados em ambos os meses
    valid_rows = (pivot[month1] > 0) & (pivot[month2] > 0)
    variation_data = pivot[valid_rows].reset_index()
    
    if variation_data.empty:
        st.warning("Não há dados suficientes para calcular a variação percentual.")
        return None
    
    # Criar gráfico de barras para variação percentual
    fig = px.bar(
        variation_data,
        x=column,
        y='Variação',
        title=f"Variação Percentual de {column} entre {month1} e {month2}",
        color='Variação',
        color_continuous_scale=['green', 'yellow', 'red'],
        range_color=[-50, 50],
        height=600  # Aumentar altura do gráfico
    )
    
    fig.update_layout(
        xaxis_title=column,
        yaxis_title="Variação Percentual (%)",
        template="plotly_white",
        title_font_size=24,  # Aumentar tamanho do título
        font=dict(size=16),  # Aumentar tamanho da fonte geral
        margin=dict(l=50, r=50, t=80, b=50)  # Ajustar margens
    )
    
    return fig

# Função para criar gráfico de pizza
def create_pie_chart(df, column, title):
    if df.empty:
        st.warning("Não há dados para exibir no gráfico.")
        return None
    
    count_df = df[column].value_counts().reset_index()
    count_df.columns = [column, 'Contagem']
    
    fig = px.pie(
        count_df, 
        names=column, 
        values='Contagem',
        title=title,
        color_discrete_sequence=px.colors.qualitative.Set3,
        height=600  # Aumentar altura do gráfico
    )
    
    fig.update_layout(
        template="plotly_white",
        title_font_size=24,  # Aumentar tamanho do título
        font=dict(size=16),  # Aumentar tamanho da fonte geral
        legend=dict(font=dict(size=14)),  # Aumentar tamanho da fonte da legenda
        margin=dict(l=50, r=50, t=80, b=50)  # Ajustar margens
    )
    
    return fig

# Função para criar análise por tipo de crime ao longo dos meses
def create_crime_analysis(df, selected_crimes=None):
    if df.empty:
        st.warning("Não há dados para exibir no gráfico.")
        return None
    
    # Se não houver mês de referência, usar o mês da data
    if 'MES_REFERENCIA' not in df.columns:
        df['MES_REFERENCIA'] = df['DATA_HORA'].dt.month.map({
            1: 'Janeiro', 2: 'Fevereiro', 3: 'Março', 4: 'Abril',
            5: 'Maio', 6: 'Junho', 7: 'Julho', 8: 'Agosto',
            9: 'Setembro', 10: 'Outubro', 11: 'Novembro', 12: 'Dezembro'
        })
    
    # Filtrar por crimes selecionados, se houver
    if selected_crimes and len(selected_crimes) > 0:
        df = df[df['EVENTO'].isin(selected_crimes)]
    else:
        # Se não houver crimes selecionados, usar os 5 mais comuns
        top_crimes = df['EVENTO'].value_counts().nlargest(5).index.tolist()
        df = df[df['EVENTO'].isin(top_crimes)]
    
    # Agrupar por mês e tipo de crime
    grouped = df.groupby(['MES_REFERENCIA', 'EVENTO']).size().reset_index(name='Contagem')
    
    # Ordenar os meses corretamente
    month_order = {month: i for i, month in enumerate(MESES)}
    grouped['month_order'] = grouped['MES_REFERENCIA'].map(month_order)
    grouped = grouped.sort_values('month_order')
    
    # Criar gráfico de linhas
    fig = px.line(
        grouped,
        x='MES_REFERENCIA',
        y='Contagem',
        color='EVENTO',
        title="Análise de Crimes por Mês",
        markers=True,
        height=600,  # Aumentar altura do gráfico
        category_orders={"MES_REFERENCIA": MESES}  # Garantir ordem correta dos meses
    )
    
    fig.update_layout(
        xaxis_title="Mês",
        yaxis_title="Número de Ocorrências",
        legend_title="Tipo de Crime",
        template="plotly_white",
        title_font_size=24,  # Aumentar tamanho do título
        font=dict(size=16),  # Aumentar tamanho da fonte geral
        legend=dict(font=dict(size=14)),  # Aumentar tamanho da fonte da legenda
        margin=dict(l=50, r=50, t=80, b=50)  # Ajustar margens
    )
    
    # Aumentar tamanho dos marcadores e linhas
    fig.update_traces(
        marker=dict(size=12),
        line=dict(width=3)
    )
    
    return fig

# Função para geocodificar endereços
@st.cache_data
def geocode_address(municipio, logradouro, numero, bairro):
    try:
        # Inicializar o geocodificador
        geolocator = Nominatim(user_agent="crime_analysis_app")
        
        # Construir o endereço completo
        address = f"{logradouro}, {numero}, {bairro}, {municipio}, Brasil"
        
        # Geocodificar o endereço
        location = geolocator.geocode(address, timeout=10)
        
        # Se não encontrar, tentar sem o número
        if location is None:
            address = f"{logradouro}, {bairro}, {municipio}, Brasil"
            location = geolocator.geocode(address, timeout=10)
        
        # Se ainda não encontrar, tentar apenas com município e bairro
        if location is None:
            address = f"{bairro}, {municipio}, Brasil"
            location = geolocator.geocode(address, timeout=10)
            
        # Se ainda não encontrar, tentar apenas com município
        if location is None:
            address = f"{municipio}, Brasil"
            location = geolocator.geocode(address, timeout=10)
        
        if location:
            return (location.latitude, location.longitude)
        else:
            return None
    except (GeocoderTimedOut, GeocoderUnavailable):
        # Em caso de timeout ou serviço indisponível
        return None
    except Exception as e:
        st.error(f"Erro ao geocodificar endereço: {e}")
        return None

# Função para criar mapa de calor usando endereços
def create_heatmap_from_addresses(df):
    if df.empty:
        st.warning("Não há dados para exibir no mapa.")
        return None
    
    # Verificar se as colunas necessárias existem
    required_columns = ['MUNICÍPIO', 'LOGRADOURO', 'NÚMERO DO LOGRADOURO', 'BAIRRO']
    missing_columns = [col for col in required_columns if col not in df.columns]
    
    if missing_columns:
        st.warning(f"Colunas necessárias ausentes: {', '.join(missing_columns)}")
        return None
    
    # Criar um DataFrame para armazenar as coordenadas
    coords_df = pd.DataFrame(columns=['lat', 'lon'])
    
    # Barra de progresso
    progress_bar = st.progress(0)
    status_text = st.empty()
    
    # Limitar a quantidade de endereços para geocodificação (para evitar sobrecarga da API)
    sample_size = min(500, len(df))
    df_sample = df.sample(sample_size) if len(df) > sample_size else df
    
    # Geocodificar endereços
    status_text.text("Geocodificando endereços... Isso pode levar alguns minutos.")
    
    coords_list = []
    for i, (_, row) in enumerate(df_sample.iterrows()):
        # Atualizar barra de progresso
        progress = int((i + 1) / len(df_sample) * 100)
        progress_bar.progress(progress)
        
        # Obter valores das colunas, tratando valores ausentes
        municipio = str(row['MUNICÍPIO']) if pd.notna(row['MUNICÍPIO']) else ""
        logradouro = str(row['LOGRADOURO']) if pd.notna(row['LOGRADOURO']) else ""
        numero = str(row['NÚMERO DO LOGRADOURO']) if pd.notna(row['NÚMERO DO LOGRADOURO']) else ""
        bairro = str(row['BAIRRO']) if pd.notna(row['BAIRRO']) else ""
        
        # Pular se não houver informações suficientes
        if not (municipio and (logradouro or bairro)):
            continue
        
        # Geocodificar o endereço
        coords = geocode_address(municipio, logradouro, numero, bairro)
        
        if coords:
            coords_list.append(coords)
        
        # Adicionar um pequeno atraso para evitar sobrecarregar a API
        time.sleep(0.1)
    
    # Limpar a barra de progresso e o texto de status
    progress_bar.empty()
    status_text.empty()
    
    # Verificar se há coordenadas válidas
    if not coords_list:
        st.warning("Não foi possível geocodificar nenhum endereço. Verifique os dados de endereço.")
        return None
    
    # Criar mapa centrado na média das coordenadas
    center_lat = sum(lat for lat, _ in coords_list) / len(coords_list)
    center_lon = sum(lon for _, lon in coords_list) / len(coords_list)
    
    m = folium.Map(location=[center_lat, center_lon], zoom_start=12, width='100%')
    
    # Adicionar pontos de calor
    HeatMap(coords_list).add_to(m)
    
    return m

# Função para criar mapa de calor usando coordenadas existentes
def create_heatmap_from_coordinates(df):
    if df.empty:
        st.warning("Não há dados para exibir no mapa.")
        return None
    
    # Verificar se há coordenadas válidas
    df['COORDENADA X'] = pd.to_numeric(df['COORDENADA X'], errors='coerce')
    df['COORDENADA y'] = pd.to_numeric(df['COORDENADA y'], errors='coerce')
    
    # Filtrar apenas registros com coordenadas válidas
    valid_coords = df.dropna(subset=['COORDENADA X', 'COORDENADA y'])
    
    if valid_coords.empty:
        st.warning("Não há coordenadas válidas para exibir no mapa.")
        return None
    
    # Criar mapa centrado na média das coordenadas
    center_lat = valid_coords['COORDENADA y'].mean()
    center_lon = valid_coords['COORDENADA X'].mean()
    
    m = folium.Map(location=[center_lat, center_lon], zoom_start=12, width='100%')
    
    # Adicionar pontos de calor
    heat_data = [[row['COORDENADA y'], row['COORDENADA X']] for _, row in valid_coords.iterrows()]
    HeatMap(heat_data).add_to(m)
    
    return m

# Função para exportar para Excel
def export_to_excel(df):
    output = io.BytesIO()
    
    # Criar um novo workbook
    wb = Workbook()
    ws = wb.active
    ws.title = "Dados Criminais"
    
    # Adicionar os dados
    for r in dataframe_to_rows(df, index=False, header=True):
        ws.append(r)
    
    # Criar uma nova aba para gráficos
    ws_charts = wb.create_sheet(title="Gráficos")
    
    # Adicionar dados para gráficos
    crime_counts = df['EVENTO'].value_counts().reset_index()
    crime_counts.columns = ['Tipo de Crime', 'Contagem']
    
    for r in dataframe_to_rows(crime_counts, index=False, header=True):
        ws_charts.append(r)
    
    # Criar gráfico de barras
    chart = BarChart()
    chart.title = "Ocorrências por Tipo de Crime"
    chart.x_axis.title = "Tipo de Crime"
    chart.y_axis.title = "Contagem"
    
    data = Reference(ws_charts, min_col=2, min_row=1, max_row=len(crime_counts)+1)
    cats = Reference(ws_charts, min_col=1, min_row=2, max_row=len(crime_counts)+1)
    chart.add_data(data, titles_from_data=True)
    chart.set_categories(cats)
    
    ws_charts.add_chart(chart, "A10")
    
    # Salvar o workbook
    wb.save(output)
    output.seek(0)
    
    return output

# Função para exportar para PowerPoint
def export_to_ppt(df, bar_fig, pie_fig, analysis_fig, comparative_fig=None):
    # Criar apresentação
    prs = Presentation()
    
    # Slide de título
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Análise de Dados Criminais"
    
    # Verificar se há múltiplos meses
    if 'MES_REFERENCIA' in df.columns and df['MES_REFERENCIA'].nunique() > 1:
        months = sorted(df['MES_REFERENCIA'].unique())
        subtitle.text = f"Análise Comparativa: {', '.join(months)}"
    else:
        subtitle.text = f"Período: {df['DATA_HORA'].min().strftime('%d/%m/%Y')} a {df['DATA_HORA'].max().strftime('%d/%m/%Y')}"
    
    # Salvar gráficos como imagens temporárias
    temp_files = []
    
    if bar_fig:
        bar_img_path = f"temp_bar_chart_{uuid.uuid4()}.png"
        bar_fig.write_image(bar_img_path)
        temp_files.append(bar_img_path)
        
        # Slide para gráfico de barras
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide.shapes.title
        title.text = "Ocorrências por Tipo de Crime"
        
        slide.shapes.add_picture(bar_img_path, Inches(1), Inches(1.5), width=Inches(8))
    
    if pie_fig:
        pie_img_path = f"temp_pie_chart_{uuid.uuid4()}.png"
        pie_fig.write_image(pie_img_path)
        temp_files.append(pie_img_path)
        
        # Slide para gráfico de pizza
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide.shapes.title
        title.text = "Proporção por Tipo de Crime"
        
        slide.shapes.add_picture(pie_img_path, Inches(1), Inches(1.5), width=Inches(8))
    
    if analysis_fig:
        analysis_img_path = f"temp_analysis_chart_{uuid.uuid4()}.png"
        analysis_fig.write_image(analysis_img_path)
        temp_files.append(analysis_img_path)
        
        # Slide para análise de crimes
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide.shapes.title
        title.text = "Análise de Crimes por Mês"
        
        slide.shapes.add_picture(analysis_img_path, Inches(1), Inches(1.5), width=Inches(8))
    
    # Adicionar gráfico comparativo se disponível
    if comparative_fig:
        comp_img_path = f"temp_comparative_chart_{uuid.uuid4()}.png"
        comparative_fig.write_image(comp_img_path)
        temp_files.append(comp_img_path)
        
        # Slide para gráfico comparativo
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide.shapes.title
        title.text = "Análise Comparativa entre Períodos"
        
        slide.shapes.add_picture(comp_img_path, Inches(1), Inches(1.5), width=Inches(8))
    
    # Slide para tabela de totalizadores
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "Resumo Estatístico"
    
    # Criar tabela de totalizadores
    crime_counts = df['EVENTO'].value_counts().reset_index()
    crime_counts.columns = ['Tipo de Crime', 'Contagem']
    
    rows, cols = len(crime_counts) + 1, 2
    left = Inches(2)
    top = Inches(2)
    width = Inches(6)
    height = Inches(0.5 * rows)
    
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # Preencher cabeçalho
    table.cell(0, 0).text = "Tipo de Crime"
    table.cell(0, 1).text = "Contagem"
    
    # Preencher dados
    for i, (crime, count) in enumerate(zip(crime_counts['Tipo de Crime'], crime_counts['Contagem'])):
        table.cell(i+1, 0).text = str(crime)
        table.cell(i+1, 1).text = str(count)
    
    # Slide para análise textual
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "Análise Textual"
    
    # Gerar análise textual simples
    total_ocorrencias = len(df)
    crime_mais_comum = crime_counts.iloc[0]['Tipo de Crime']
    qtd_crime_mais_comum = crime_counts.iloc[0]['Contagem']
    
    text_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4))
    text_frame = text_box.text_frame
    
    p = text_frame.add_paragraph()
    p.text = f"Análise de {total_ocorrencias} ocorrências registradas no período."
    p.font.size = Pt(14)
    
    p = text_frame.add_paragraph()
    p.text = f"O tipo de crime mais comum foi '{crime_mais_comum}' com {qtd_crime_mais_comum} ocorrências, representando {(qtd_crime_mais_comum/total_ocorrencias*100):.1f}% do total."
    p.font.size = Pt(14)
    
    # Adicionar análise comparativa se houver múltiplos meses
    if 'MES_REFERENCIA' in df.columns and df['MES_REFERENCIA'].nunique() > 1:
        p = text_frame.add_paragraph()
        p.text = f"A análise comparativa entre {df['MES_REFERENCIA'].nunique()} períodos mostra variações nos padrões criminais ao longo do tempo."
        p.font.size = Pt(14)
        
        # Análise por mês
        monthly_data = df.groupby('MES_REFERENCIA').size()
        max_month = monthly_data.idxmax()
        min_month = monthly_data.idxmin()
        
        p = text_frame.add_paragraph()
        p.text = f"O período com maior número de ocorrências foi {max_month} com {monthly_data[max_month]} registros."
        p.font.size = Pt(14)
        
        p = text_frame.add_paragraph()
        p.text = f"O período com menor número de ocorrências foi {min_month} com {monthly_data[min_month]} registros."
        p.font.size = Pt(14)
    else:
        p = text_frame.add_paragraph()
        p.text = f"A média mensal de ocorrências no período analisado foi de {total_ocorrencias / df['DATA_HORA'].dt.to_period('M').nunique():.1f} registros."
        p.font.size = Pt(14)
    
    # Adicionar créditos do autor
    p = text_frame.add_paragraph()
    p.text = "Criado por Leandro Vieira de Souza"
    p.font.size = Pt(12)
    p.font.italic = True
    
    # Salvar apresentação em um buffer
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    
    # Limpar arquivos temporários
    for file in temp_files:
        if os.path.exists(file):
            os.remove(file)
    
    return output

# Função para obter todas as unidades únicas do DataFrame
def get_unique_units(df):
    all_units = []
    
    # Iterar sobre todas as linhas e extrair unidades
    for unit_string in df['UNIDADE DA VIATURA'].dropna():
        units = extract_units(unit_string)
        all_units.extend(units)
    
    # Remover duplicatas e ordenar
    unique_units = sorted(list(set(all_units)))
    return unique_units

# Interface principal
def main():
    # Título e descrição
    with st.container():
        col1, col2 = st.columns([4, 1])
        with col1:
            st.title("🚨 Análise de Dados Criminais - PM")
            st.markdown("Sistema de análise de ocorrências registradas pela Polícia Militar")
    
    # Criar layout com colunas para maximizar espaço
    col_sidebar, col_main = st.columns([1, 4])
    
    # Sidebar para upload e filtros
    with col_sidebar:
        with st.expander("📤 Upload de Dados", expanded=True):
            # Opção para upload de múltiplas planilhas
            upload_option = st.radio(
                "Escolha o modo de upload:",
                ["Upload de planilha única", "Upload de múltiplas planilhas (comparação mensal)"]
            )
            
            if upload_option == "Upload de planilha única":
                uploaded_file = st.file_uploader("Carregar planilha de ocorrências", type=["xlsx"])
                
                if uploaded_file:
                    # Selecionar o mês de referência
                    month_name = st.selectbox("Selecione o mês de referência:", MESES)
                    
                    # Carregar dados
                    df = load_data(uploaded_file, month_name)
                    
                    # Armazenar no estado da sessão
                    st.session_state.dataframes[month_name] = df
                    st.session_state.active_dataframes = [month_name]
                    
                    st.success(f"Dados de {month_name} carregados com sucesso! {len(df)} registros encontrados.")
            else:
                # Upload de múltiplas planilhas
                st.markdown("### Upload de Planilhas Mensais")
                st.info("Faça upload de planilhas de diferentes meses para comparação.")
                
                # Área para upload de múltiplas planilhas
                uploaded_file = st.file_uploader(
                    "Carregar planilha mensal", 
                    type=["xlsx"]
                )
                
                if uploaded_file:
                    # Selecionar o mês de referência
                    month_name = st.selectbox("Selecione o mês de referência:", MESES)
                    
                    # Botão para adicionar a planilha
                    if st.button("Adicionar Planilha"):
                        # Verificar se o mês já foi carregado
                        if month_name in st.session_state.dataframes:
                            st.warning(f"Já existe uma planilha para {month_name}. Ela será substituída.")
                        
                        # Carregar dados
                        df = load_data(uploaded_file, month_name)
                        
                        # Armazenar no estado da sessão
                        st.session_state.dataframes[month_name] = df
                        
                        # Adicionar à lista de ativos se não estiver lá
                        if month_name not in st.session_state.active_dataframes:
                            st.session_state.active_dataframes.append(month_name)
                        
                        st.success(f"Planilha de {month_name} adicionada com sucesso! {len(df)} registros.")
                
                # Mostrar quais planilhas foram carregadas
                if st.session_state.dataframes:
                    st.markdown("### Planilhas Carregadas")
                    
                    for month, df in st.session_state.dataframes.items():
                        st.info(f"{month}: {len(df)} registros")
        
        # Verificar se há dados para mostrar filtros
        if st.session_state.dataframes and st.session_state.active_dataframes:
            with st.expander("🔍 Filtros", expanded=True):
                # Seleção de meses para análise
                st.subheader("Meses para Análise")
                all_months = list(st.session_state.dataframes.keys())
                selected_months = st.multiselect(
                    "Selecione os meses para incluir na análise:",
                    all_months,
                    default=st.session_state.active_dataframes
                )
                
                # Atualizar meses ativos
                if selected_months:
                    st.session_state.active_dataframes = selected_months
                
                # Combinar os DataFrames ativos
                df = combine_dataframes(st.session_state.dataframes, st.session_state.active_dataframes)
                
                # Filtro de data
                st.subheader("Período")
                min_date = df['DATA_HORA'].min().date()
                max_date = df['DATA_HORA'].max().date()
                
                col1, col2 = st.columns(2)
                with col1:
                    start_date = st.date_input("Data inicial", min_date, format="DD/MM/YYYY")
                with col2:
                    end_date = st.date_input("Data final", max_date, format="DD/MM/YYYY")
                
                # Filtro de tipo de crime
                st.subheader("Tipo de Crime")
                crime_options = sorted(df['EVENTO'].unique())
                crime_type = st.multiselect("Selecione os tipos de crime", crime_options)
                
                # Filtro de localidade
                st.subheader("Localidade")
                location_options = sorted(df['ÁREA URBANA'].dropna().unique())
                location = st.multiselect("Selecione as localidades", location_options)
                
                # Filtro de unidade responsável - modificado para mostrar unidades individuais
                st.subheader("Unidade Responsável")
                unit_options = get_unique_units(df)  # Obter unidades únicas
                unit = st.multiselect("Selecione as unidades", unit_options)
                
                # Filtro de palavras-chave
                st.subheader("Palavras-chave")
                keywords = st.text_input("Buscar nos históricos e evoluções")
                
                # Aplicar filtros
                filtered_df = filter_data(df, start_date, end_date, crime_type, location, unit, keywords)
                
                st.info(f"Exibindo {len(filtered_df)} de {len(df)} registros após aplicação dos filtros.")
            
            # Botões de exportação
            with st.expander("📊 Exportar Resultados", expanded=True):
                col1, col2 = st.columns(2)
                
                with col1:
                    if st.button("📥 Exportar Excel", use_container_width=True):
                        excel_data = export_to_excel(filtered_df)
                        st.download_button(
                            label="Baixar arquivo Excel",
                            data=excel_data,
                            file_name="dados_criminais.xlsx",
                            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                            use_container_width=True
                        )
                
                with col2:
                    if st.button("📊 Exportar PowerPoint", use_container_width=True):
                        # Verificar se estamos em modo de comparação
                        is_comparison_mode = 'MES_REFERENCIA' in filtered_df.columns and filtered_df['MES_REFERENCIA'].nunique() > 1
                        
                        # Criar gráficos para o PowerPoint
                        bar_fig = create_bar_chart(filtered_df, 'EVENTO', "Ocorrências por Tipo de Crime")
                        pie_fig = create_pie_chart(filtered_df, 'EVENTO', "Proporção por Tipo de Crime")
                        analysis_fig = create_crime_analysis(filtered_df)
                        
                        if is_comparison_mode:
                            comp_fig = create_comparative_bar_chart(filtered_df, 'EVENTO')
                            ppt_data = export_to_ppt(filtered_df, bar_fig, pie_fig, analysis_fig, comp_fig)
                        else:
                            ppt_data = export_to_ppt(filtered_df, bar_fig, pie_fig, analysis_fig)
                        
                        st.download_button(
                            label="Baixar apresentação PowerPoint",
                            data=ppt_data,
                            file_name="analise_criminal.pptx",
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                            use_container_width=True
                        )
                
                # Salvar estado dos filtros
                if st.button("💾 Salvar Filtros", use_container_width=True):
                    st.session_state.saved_filters = {
                        'start_date': start_date,
                        'end_date': end_date,
                        'crime_type': crime_type,
                        'location': location,
                        'unit': unit,
                        'keywords': keywords
                    }
                    st.success("Filtros salvos com sucesso!")
                
                # Carregar filtros salvos
                if st.button("📂 Carregar Filtros Salvos", use_container_width=True) and hasattr(st.session_state, 'saved_filters'):
                    st.write("Filtros carregados:")
                    st.write(st.session_state.saved_filters)
    
    # Conteúdo principal
    with col_main:
        if st.session_state.dataframes and st.session_state.active_dataframes:
            # Combinar os DataFrames ativos
            df = combine_dataframes(st.session_state.dataframes, st.session_state.active_dataframes)
            
            # Aplicar filtros
            filtered_df = filter_data(df, start_date, end_date, crime_type, location, unit, keywords) if 'start_date' in locals() else df
            
            if not filtered_df.empty:
                # Métricas principais
                st.header("📊 Métricas Principais")
                
                col1, col2, col3 = st.columns(3)
                
                with col1:
                    st.metric("Total de Ocorrências", len(filtered_df))
                
                with col2:
                    top_crime = filtered_df['EVENTO'].value_counts().index[0] if not filtered_df.empty else "N/A"
                    st.metric("Crime Mais Comum", top_crime)
                
                with col3:
                    top_location = filtered_df['ÁREA URBANA'].value_counts().index[0] if not filtered_df.empty else "N/A"
                    st.metric("Localidade Mais Afetada", top_location)
                
                # Visualizações
                st.header("📈 Visualizações")
                
                # Seletor de mês específico para visualizações
                if len(st.session_state.active_dataframes) > 1:
                    selected_month_viz = st.selectbox(
                        "Selecione um mês específico para visualização:",
                        ["Todos os meses selecionados"] + st.session_state.active_dataframes
                    )
                    
                    if selected_month_viz != "Todos os meses selecionados":
                        # Filtrar apenas o mês selecionado
                        viz_df = filtered_df[filtered_df['MES_REFERENCIA'] == selected_month_viz]
                    else:
                        viz_df = filtered_df
                else:
                    viz_df = filtered_df
                
                tab1, tab2, tab3, tab4 = st.tabs([
                    "Gráficos de Barras", 
                    "Gráficos de Pizza", 
                    "Análise", 
                    "Mapa de Calor"
                ])
                
                with tab1:
                    st.subheader("Ocorrências por Tipo de Crime")
                    bar_fig = create_bar_chart(viz_df, 'EVENTO', "Ocorrências por Tipo de Crime")
                    if bar_fig:
                        st.plotly_chart(bar_fig, use_container_width=True)
                    
                    st.subheader("Ocorrências por Localidade")
                    bar_fig_loc = create_bar_chart(viz_df, 'ÁREA URBANA', "Ocorrências por Localidade", color='#15803D')
                    if bar_fig_loc:
                        st.plotly_chart(bar_fig_loc, use_container_width=True)
                
                with tab2:
                    st.subheader("Proporção por Tipo de Crime")
                    pie_fig = create_pie_chart(viz_df, 'EVENTO', "Proporção por Tipo de Crime")
                    if pie_fig:
                        st.plotly_chart(pie_fig, use_container_width=True)
                
                with tab3:
                    st.subheader("Análise de Crimes por Mês")
                    
                    # Seleção de crimes para análise
                    crime_options = sorted(viz_df['EVENTO'].unique())
                    selected_crimes = st.multiselect(
                        "Selecione os tipos de crime para analisar:",
                        crime_options,
                        default=viz_df['EVENTO'].value_counts().nlargest(5).index.tolist()
                    )
                    
                    # Criar gráfico de análise
                    analysis_fig = create_crime_analysis(viz_df, selected_crimes)
                    if analysis_fig:
                        st.plotly_chart(analysis_fig, use_container_width=True)
                    
                    # Adicionar explicação
                    st.markdown("""
                    <div style="background-color: #f0f2f6; padding: 1rem; border-radius: 0.5rem; margin-top: 1rem;">
                        <h4 style="margin-top: 0;">Sobre esta Análise</h4>
                        <p>
                            Este gráfico mostra a evolução de cada tipo de crime ao longo dos meses selecionados.
                            Cada linha colorida representa um tipo específico de crime, permitindo visualizar
                            tendências, sazonalidades e comparar a incidência de diferentes crimes no mesmo período.
                        </p>
                    </div>
                    """, unsafe_allow_html=True)
                
                with tab4:
                    st.subheader("Mapa de Calor de Ocorrências")
                    
                    # Opções para o mapa de calor
                    map_option = st.radio(
                        "Escolha o método para gerar o mapa de calor:",
                        ["Usar endereços (MUNICÍPIO, LOGRADOURO, BAIRRO)", "Usar coordenadas (X, Y)"]
                    )
                    
                    if map_option == "Usar endereços (MUNICÍPIO, LOGRADOURO, BAIRRO)":
                        heatmap = create_heatmap_from_addresses(viz_df)
                    else:
                        heatmap = create_heatmap_from_coordinates(viz_df)
                    
                    if heatmap:
                        # Aumentar tamanho do mapa
                        folium_static(heatmap, width=1200, height=700)
                    else:
                        st.warning("Não foi possível gerar o mapa de calor. Verifique se há dados de localização válidos.")
                
                # Seção autônoma para análise comparativa
                if len(st.session_state.active_dataframes) > 1:
                    st.header("🔄 Análise Criminal Comparativa")
                    
                    st.markdown("""
                    <div style="background-color: #f0f2f6; padding: 1rem; border-radius: 0.5rem; margin-bottom: 1rem;">
                        <h3 style="margin-top: 0;">Comparação de Índices Criminais</h3>
                        <p style="font-size: 1.1rem;">
                            Esta seção permite comparar índices criminais entre diferentes meses para 
                            identificar tendências, aumentos ou diminuições nos tipos de crimes.
                        </p>
                    </div>
                    """, unsafe_allow_html=True)
                    
                    # Seleção de meses para comparação
                    comp_months = st.multiselect(
                        "Selecione os meses para comparação:",
                        st.session_state.active_dataframes,
                        default=st.session_state.active_dataframes[:min(2, len(st.session_state.active_dataframes))]
                    )
                    
                    if len(comp_months) >= 2:
                        # Filtrar dados apenas para os meses selecionados
                        comp_df = filtered_df[filtered_df['MES_REFERENCIA'].isin(comp_months)]
                        
                        # Seleção de tipos de crime para comparação
                        selected_crimes = st.multiselect(
                            "Selecione os tipos de crime para comparar:",
                            sorted(comp_df['EVENTO'].unique()),
                            default=comp_df['EVENTO'].value_counts().nlargest(5).index.tolist(),
                            key="comp_crimes"
                        )
                        
                        if selected_crimes:
                            # Gráfico de barras comparativo
                            st.subheader("Comparação de Crimes por Mês")
                            comp_df_filtered = comp_df[comp_df['EVENTO'].isin(selected_crimes)]
                            comp_bar_fig = create_comparative_bar_chart(comp_df_filtered, 'EVENTO')
                            if comp_bar_fig:
                                st.plotly_chart(comp_bar_fig, use_container_width=True)
                            
                            # Análise de variação percentual
                            if len(comp_months) == 2:
                                st.subheader("Variação Percentual entre Períodos")
                                var_fig = create_percentage_change_chart(comp_df_filtered, 'EVENTO', comp_months)
                                if var_fig:
                                    st.plotly_chart(var_fig, use_container_width=True)
                                    
                                    # Calcular estatísticas de variação
                                    grouped = comp_df_filtered.groupby(['MES_REFERENCIA', 'EVENTO']).size().reset_index(name='Contagem')
                                    pivot = grouped.pivot(index='EVENTO', columns='MES_REFERENCIA', values='Contagem').fillna(0)
                                    
                                    month1, month2 = comp_months[0], comp_months[1]
                                    if month1 in pivot.columns and month2 in pivot.columns:
                                        pivot['Variação'] = ((pivot[month2] - pivot[month1]) / pivot[month1] * 100).fillna(0)
                                        
                                        # Filtrar apenas os tipos de crime com dados em ambos os meses
                                        valid_rows = (pivot[month1] > 0) & (pivot[month2] > 0)
                                        variation_data = pivot[valid_rows]
                                        
                                        if not variation_data.empty:
                                            # Calcular estatísticas
                                            aumentos = (variation_data['Variação'] > 0).sum()
                                            diminuicoes = (variation_data['Variação'] < 0).sum()
                                            sem_alteracao = (variation_data['Variação'] == 0).sum()
                                            
                                            # Mostrar estatísticas em cards
                                            st.markdown("""
                                            <h3 style="margin-top: 1.5rem;">Resumo da Variação</h3>
                                            """, unsafe_allow_html=True)
                                            
                                            col1, col2, col3 = st.columns(3)
                                            
                                            with col1:
                                                st.markdown(f"""
                                                <div style="background-color: #ffcccb; padding: 1rem; border-radius: 0.5rem; text-align: center;">
                                                    <h2 style="margin: 0; color: #cc0000;">{aumentos}</h2>
                                                    <p style="margin: 0; font-weight: bold;">Crimes com Aumento</p>
                                                </div>
                                                """, unsafe_allow_html=True)
                                            
                                            with col2:
                                                st.markdown(f"""
                                                <div style="background-color: #ccffcc; padding: 1rem; border-radius: 0.5rem; text-align: center;">
                                                    <h2 style="margin: 0; color: #007700;">{diminuicoes}</h2>
                                                    <p style="margin: 0; font-weight: bold;">Crimes com Diminuição</p>
                                                </div>
                                                """, unsafe_allow_html=True)
                                            
                                            with col3:
                                                st.markdown(f"""
                                                <div style="background-color: #e0e0e0; padding: 1rem; border-radius: 0.5rem; text-align: center;">
                                                    <h2 style="margin: 0; color: #555555;">{sem_alteracao}</h2>
                                                    <p style="margin: 0; font-weight: bold;">Sem Alteração</p>
                                                </div>
                                                """, unsafe_allow_html=True)
                                            
                                            # Mostrar os maiores aumentos e diminuições
                                            col1, col2 = st.columns(2)
                                            
                                            with col1:
                                                st.markdown("""
                                                <h4 style="margin-top: 1.5rem;">Maiores Aumentos:</h4>
                                                """, unsafe_allow_html=True)
                                                
                                                # Mostrar os 3 maiores aumentos
                                                top_increases = variation_data.sort_values('Variação', ascending=False).head(3)
                                                for crime, row in top_increases.iterrows():
                                                    st.markdown(f"""
                                                    <div style="background-color: #fff0f0; padding: 0.8rem; border-radius: 0.5rem; margin-bottom: 0.5rem;">
                                                        <h5 style="margin: 0; color: #cc0000;">{crime}</h5>
                                                        <p style="margin: 0; font-weight: bold;">Aumento de {row['Variação']:.1f}%</p>
                                                        <p style="margin: 0;">({int(row[month1])} → {int(row[month2])} ocorrências)</p>
                                                    </div>
                                                    """, unsafe_allow_html=True)
                                            
                                            with col2:
                                                st.markdown("""
                                                <h4 style="margin-top: 1.5rem;">Maiores Diminuições:</h4>
                                                """, unsafe_allow_html=True)
                                                
                                                # Mostrar as 3 maiores diminuições
                                                top_decreases = variation_data.sort_values('Variação').head(3)
                                                for crime, row in top_decreases.iterrows():
                                                    st.markdown(f"""
                                                    <div style="background-color: #f0fff0; padding: 0.8rem; border-radius: 0.5rem; margin-bottom: 0.5rem;">
                                                        <h5 style="margin: 0; color: #007700;">{crime}</h5>
                                                        <p style="margin: 0; font-weight: bold;">Diminuição de {abs(row['Variação']):.1f}%</p>
                                                        <p style="margin: 0;">({int(row[month1])} → {int(row[month2])} ocorrências)</p>
                                                    </div>
                                                    """, unsafe_allow_html=True)
                            
                            # Tabela comparativa
                            st.subheader("Tabela Comparativa por Mês")
                            
                            # Criar tabela pivô
                            pivot_table = pd.pivot_table(
                                comp_df_filtered,
                                values='ID',
                                index=['EVENTO'],
                                columns=['MES_REFERENCIA'],
                                aggfunc='count',
                                fill_value=0
                            )
                            
                            # Adicionar linha de total
                            pivot_table.loc['TOTAL'] = pivot_table.sum()
                            
                            # Estilizar a tabela para melhor visualização
                            st.dataframe(
                                pivot_table,
                                use_container_width=True,
                                height=400
                            )
                        else:
                            st.warning("Selecione pelo menos um tipo de crime para comparação.")
                    else:
                        st.warning("Selecione pelo menos dois meses para comparação.")
            
            else:
                st.warning("Nenhum dado encontrado com os filtros aplicados. Tente ajustar os critérios de filtro.")
        
        else:
            st.info("Por favor, faça o upload de uma ou mais planilhas Excel (.xlsx) para começar a análise.")
            
            # Exemplo de como os dados devem estar estruturados
            with st.expander("ℹ️ Informações sobre o formato da planilha", expanded=False):
                st.markdown("""
                    A planilha deve conter as seguintes colunas:
                    - DATA DE INÍCIO DO ATENDIMENTO (formato DD/MM/AAAA)
                    - HORA DE INÍCIO DO ATENDIMENTO
                    - UNIDADE DA VIATURA
                    - EVENTO
                    - CIRCUNSTÂNCIA
                    - ÁREA URBANA
                    - MUNICÍPIO
                    - LOGRADOURO
                    - NÚMERO DO LOGRADOURO
                    - BAIRRO
                    - HISTÓRICOS
                    - EVOLUÇÕES
                """)

if __name__ == "__main__":
    main()
